# -*- coding: utf-8 -*-
"""MedVisual sunum (.pptx) olusturucu — 16:9, Turkce, mobil + web ekran goruntuleri."""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from PIL import Image as PILImage

HERE = os.path.dirname(os.path.abspath(__file__))
MOBIL = os.path.join(HERE, "mobil")
WEB = os.path.join(HERE, "web")
OUT = os.path.abspath(os.path.join(HERE, "..", "MedVisual_Sunum.pptx"))

# --- Renkler ---
INDIGO = RGBColor(0x3F, 0x51, 0xB5)
INDIGO_D = RGBColor(0x30, 0x3F, 0x9F)
TEAL = RGBColor(0x00, 0x89, 0x7B)
TEAL_D = RGBColor(0x00, 0x6B, 0x60)
BG = RGBColor(0xF7, 0xF8, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1F, 0x36)
SUB = RGBColor(0x5B, 0x61, 0x77)
LIGHT = RGBColor(0xE6, 0xE9, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_SLIDE = RGBColor(0x16, 0x1B, 0x33)
SHADOW = RGBColor(0xD7, 0xDB, 0xE8)

EMU_IN = 914400
SW, SH = 13.333, 7.5  # 16:9

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


def img_exists(path):
    return os.path.isfile(path)


def img_ratio(path):
    with PILImage.open(path) as im:
        return im.width / im.height


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def grad(shape, c1, c2, angle=45):
    """Basit iki renkli gradyan."""
    from pptx.oxml.ns import qn
    sp = shape.fill._xPr  # spPr
    # mevcut fill'i temizle
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for e in sp.findall(qn(tag)):
            sp.remove(e)
    gradFill = sp.makeelement(qn("a:gradFill"), {})
    gsLst = gradFill.makeelement(qn("a:gsLst"), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = gradFill.makeelement(qn("a:gs"), {"pos": str(pos)})
        srgb = gradFill.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (col[0], col[1], col[2])})
        gs.append(srgb)
        gsLst.append(gs)
    gradFill.append(gsLst)
    lin = gradFill.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000)), "scaled": "1"})
    gradFill.append(lin)
    # line'dan once ekle
    ln = sp.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(gradFill)
    else:
        sp.append(gradFill)
    shape.line.fill.background()


def rect(slide, x, y, w, h, color=None, shape_type=MSO_SHAPE.RECTANGLE, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if color is not None:
        fill(sp, color)
    else:
        sp.fill.background()
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    else:
        sp.line.fill.background()
    return sp


def softcard(slide, x, y, w, h, color=CARD, radius=True):
    """Hafif golge efektli kart (arka golge + kart)."""
    rect(slide, x + 0.04, y + 0.06, w, h, SHADOW,
         MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE)
    c = rect(slide, x, y, w, h, color,
             MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE)
    return c


def txt(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri", line_spacing=None, space_after=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=14, color=INK, bullet_color=TEAL,
            gap=8, lead_bold=False, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = line_spacing
        rb = p.add_run()
        rb.text = "•  "
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = bullet_color
        rb.font.name = "Calibri"
        # baslik:aciklama ayrimi
        if lead_bold and "—" in it:
            head, rest = it.split("—", 1)
            r1 = p.add_run(); r1.text = head.strip() + " "
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = "Calibri"
            r2 = p.add_run(); r2.text = "— " + rest.strip()
            r2.font.size = Pt(size); r2.font.color.rgb = SUB; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def base_slide(title=None, accent=INDIGO, dark_header=False):
    s = add_slide()
    rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, BG)
    if title is not None:
        # ust accent seridi
        rect(s, 0, 0, SW, 0.16, accent)
        txt(s, 0.6, 0.34, SW - 1.2, 0.7, title, size=26, color=INK, bold=True)
        # baslik alti ince cizgi
        rect(s, 0.62, 1.02, 1.5, 0.055, accent)
        rect(s, 2.12, 1.035, 1.0, 0.03, TEAL)
    return s


def place_image(slide, path, x, y, max_w, max_h, frame=True, frame_color=LIGHT,
                center=True, caption=None, cap_color=SUB):
    """Goruntuyu cerceveyle (golgeli kart) orana gore yerlestir."""
    if not img_exists(path):
        if frame:
            softcard(slide, x, y, max_w, max_h, CARD)
        txt(slide, x, y + max_h / 2 - 0.3, max_w, 0.6,
            "(görsel mevcut değil)", size=12, color=SUB, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, italic=True)
        return
    r = img_ratio(path)  # w/h
    # alana sigdir
    w = max_w
    h = w / r
    if h > max_h:
        h = max_h
        w = h * r
    ox = x + (max_w - w) / 2 if center else x
    oy = y + (max_h - h) / 2 if center else y
    pad = 0.08
    if frame:
        softcard(slide, ox - pad, oy - pad, w + 2 * pad, h + 2 * pad, CARD)
    slide.shapes.add_picture(path, Inches(ox), Inches(oy), Inches(w), Inches(h))
    if caption:
        txt(slide, ox - pad, oy + h + pad + 0.04, w + 2 * pad, 0.4, caption,
            size=11.5, color=cap_color, align=PP_ALIGN.CENTER, italic=True)
    return (ox, oy, w, h)


def badge(slide, x, y, text, color=TEAL, w=2.2, h=0.42, fg=WHITE, size=12):
    b = rect(slide, x, y, w, h, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = fg; r.font.name = "Calibri"
    return b


def M(name):
    return os.path.join(MOBIL, name)


def W(name):
    return os.path.join(WEB, name)


# =====================================================================
# SLIDE 1 — KAPAK
# =====================================================================
def slide_cover():
    s = add_slide()
    bg = rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, INDIGO)
    grad(bg, (0x2A, 0x32, 0x80), (0x00, 0x89, 0x7B), angle=30)
    # dekoratif daireler
    c1 = rect(s, 10.2, -1.4, 4.5, 4.5, None, MSO_SHAPE.OVAL)
    fill(c1, RGBColor(0x4A, 0x5A, 0xD0)); c1.line.fill.background()
    c2 = rect(s, 9.0, 4.8, 3.2, 3.2, None, MSO_SHAPE.OVAL)
    fill(c2, RGBColor(0x12, 0x9C, 0x8C)); c2.line.fill.background()
    # logo rozet
    lg = rect(s, 0.9, 1.5, 0.95, 0.95, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = lg.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run(); r.text = "M"; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = INDIGO; r.font.name="Calibri"
    # baslik
    txt(s, 0.9, 2.7, 11.0, 1.4, "MedVisual", size=66, color=WHITE, bold=True)
    txt(s, 0.95, 3.95, 11.5, 0.7,
        "Yapay Zeka Destekli Görsel Öğrenme ve Flashcard Ekosistemi",
        size=22, color=RGBColor(0xE9,0xEC,0xFF), bold=False)
    # alt cizgi vurgu
    rect(s, 0.98, 4.62, 3.4, 0.06, RGBColor(0x6FE,0x6F,0xE0) if False else RGBColor(0x7C,0xE0,0xD4))
    txt(s, 0.95, 4.75, 11.5, 0.5,
        "Web + Mobil — Tek Ekosistem, Tek Veri Kaynağı",
        size=15.5, color=RGBColor(0xCDE9,0xCD,0xE9) if False else RGBColor(0xC9,0xF0,0xE9), bold=True)
    # footer blok
    fb = rect(s, 0.0, 6.35, SW, 1.15, None, MSO_SHAPE.RECTANGLE)
    fill(fb, RGBColor(0x10, 0x16, 0x2C)); fb.line.fill.background()
    txt(s, 0.95, 6.45, 8.5, 0.4,
        "Fırat Üniversitesi · Teknoloji Fakültesi · Yazılım Mühendisliği",
        size=13.5, color=WHITE, bold=True)
    txt(s, 0.95, 6.82, 8.5, 0.35, "215541014@firat.edu.tr",
        size=12, color=RGBColor(0x9F,0xC9,0xF0))
    txt(s, 0.95, 7.10, 11.5, 0.35,
        "Web Tasarım ve Programlama  &  Fonksiyonel Programlama",
        size=11.5, color=RGBColor(0xB8,0xBE,0xD6), italic=True)


# =====================================================================
# SLIDE 2 — PROBLEM & COZUM
# =====================================================================
def slide_problem():
    s = base_slide("Problem ve Çözüm", accent=INDIGO)
    cw = (SW - 0.6 - 0.6 - 0.5) / 2  # iki sutun
    x1 = 0.6
    x2 = 0.6 + cw + 0.5
    cy = 1.45
    ch = 5.4
    # Problem karti
    softcard(s, x1, cy, cw, ch, CARD)
    rect(s, x1, cy, cw, 0.85, RGBColor(0xFB,0xE9,0xE9), MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, x1 + 0.35, cy + 0.18, cw - 0.7, 0.5, "⚠  Problem", size=20, color=RGBColor(0xC0,0x39,0x2B), bold=True)
    bullets(s, x1 + 0.45, cy + 1.05, cw - 0.85, ch - 1.3, [
        "Tıp öğrencileri 600+ sayfalık yoğun PDF'lerle çalışıyor.",
        "Bilgi kartlarını elle hazırlamak çok yavaş ve yorucu.",
        "Anatomi büyük ölçüde görsel bir bilim; ancak görsel hafıza çoğu dijital araçta hiç kullanılmıyor.",
        "Farklı cihazlardaki notlar ve kartlar dağınık, senkron değil.",
    ], size=15, color=INK, bullet_color=RGBColor(0xC0,0x39,0x2B), gap=12, line_spacing=1.1)
    # Cozum karti
    softcard(s, x2, cy, cw, ch, CARD)
    rect(s, x2, cy, cw, 0.85, RGBColor(0xE3,0xF4,0xF1), MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, x2 + 0.35, cy + 0.18, cw - 0.7, 0.5, "✓  Çözüm", size=20, color=TEAL_D, bold=True)
    bullets(s, x2 + 0.45, cy + 1.05, cw - 0.85, ch - 1.3, [
        "PDF yükle → motor otomatik olarak görselli bilgi kartı ve quiz üretir.",
        "DIP, PDF'teki anatomik figürleri ayıklayıp doğru terimle eşler.",
        "Aralıklı tekrar (SM-2) ile en verimli zamanda çalış.",
        "Web ve mobil tek hesapta, anında senkron tek veri kaynağı.",
    ], size=15, color=INK, bullet_color=TEAL, gap=12, line_spacing=1.1)
    # ortada ok
    ar = rect(s, x1 + cw + 0.02, cy + ch/2 - 0.35, 0.46, 0.7, INDIGO, MSO_SHAPE.RIGHT_ARROW)


# =====================================================================
# SLIDE 3 — SISTEM MIMARISI
# =====================================================================
def slide_arch():
    s = base_slide("Sistem Mimarisi — Single Source of Truth", accent=TEAL)
    cx = SW / 2
    # Supabase (en ust)
    sw_, sh_ = 6.6, 1.0
    sup = softcard(s, cx - sw_/2, 1.45, sw_, sh_, INDIGO)
    grad(sup, (0x3F,0x51,0xB5),(0x30,0x3F,0x9F))
    txt(s, cx - sw_/2, 1.55, sw_, 0.45, "Supabase (Bulut)", size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx - sw_/2, 2.02, sw_, 0.35, "PostgreSQL  ·  Auth  ·  Storage", size=13, color=RGBColor(0xD7,0xDC,0xFF), align=PP_ALIGN.CENTER)
    # baglanti cizgisi
    rect(s, cx - 0.02, 2.45, 0.04, 0.45, TEAL)
    # API
    aw, ah = 5.4, 0.95
    api = softcard(s, cx - aw/2, 2.9, aw, ah, TEAL)
    grad(api, (0x00,0x9A,0x8A),(0x00,0x6B,0x60))
    txt(s, cx - aw/2, 3.0, aw, 0.4, "Merkezi API (FastAPI)", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx - aw/2, 3.42, aw, 0.32, "RESTful  ·  asenkron iş kuyruğu  ·  tek otorite", size=12, color=RGBColor(0xCDEFEA,0xCD,0xEA) if False else RGBColor(0xCF,0xEF,0xEA), align=PP_ALIGN.CENTER)
    # dagitim cizgileri
    by = 4.5
    boxes = [
        ("Web", "React (TypeScript)", INDIGO),
        ("Mobil", "Flutter (Dart)", RGBColor(0x2A,0x6F,0xDB)),
        ("DIP Motoru", "Python / OpenCV", RGBColor(0x6A,0x4C,0xA8)),
    ]
    bw, bh = 3.5, 1.25
    total = len(boxes) * bw + (len(boxes) - 1) * 0.55
    startx = cx - total / 2
    # yatay dagitim cizgisi
    rect(s, startx + bw/2, 4.0, total - bw, 0.04, TEAL)
    rect(s, cx - 0.02, 3.85, 0.04, 0.2, TEAL)
    for i, (t, sub, col) in enumerate(boxes):
        bx = startx + i * (bw + 0.55)
        rect(s, bx + bw/2 - 0.02, 4.04, 0.04, by - 4.04, TEAL)
        cardb = softcard(s, bx, by, bw, bh, CARD)
        rect(s, bx, by, 0.14, bh, col, MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, bx + 0.35, by + 0.22, bw - 0.5, 0.5, t, size=17, color=col, bold=True)
        txt(s, bx + 0.35, by + 0.72, bw - 0.5, 0.4, sub, size=12.5, color=SUB)
    # caption
    cap = rect(s, 0.6, 6.35, SW - 1.2, 0.7, RGBColor(0xEE,0xF1,0xFA), MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 0.9, 6.45, SW - 1.8, 0.5,
        "Hem web hem mobil aynı merkezi veriyi konuşur — cihazlar arası tam senkron, tek doğruluk kaynağı.",
        size=13.5, color=INDIGO_D, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 4 — TEKNOLOJI YIGINI
# =====================================================================
def slide_stack():
    s = base_slide("Teknoloji Yığını", accent=INDIGO)
    cols = [
        ("🌐  Web", INDIGO, [
            "React + Vite",
            "TypeScript",
            "Tailwind CSS",
            "TanStack Query",
            "FastAPI (REST)",
            "Supabase / PostgreSQL",
        ]),
        ("📱  Mobil", RGBColor(0x2A,0x6F,0xDB), [
            "Flutter",
            "Dart",
            "BLoC (durum yönetimi)",
            "freezed",
            "immutability",
            "saf fonksiyonlar",
        ]),
        ("🔬  Görüntü İşleme (DIP)", TEAL, [
            "Python",
            "OpenCV",
            "Tesseract OCR",
            "Otsu segmentasyonu",
            "morfolojik işlemler",
            "figür–terim eşleme",
        ]),
    ]
    cw = (SW - 0.6 - 0.6 - 2*0.45) / 3
    cy = 1.5
    ch = 5.45
    for i, (t, col, items) in enumerate(cols):
        x = 0.6 + i * (cw + 0.45)
        softcard(s, x, cy, cw, ch, CARD)
        hdr = rect(s, x, cy, cw, 0.95, col, MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        txt(s, x + 0.1, cy + 0.24, cw - 0.2, 0.5, t, size=16.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        bullets(s, x + 0.35, cy + 1.15, cw - 0.6, ch - 1.4, items, size=14.5, color=INK,
                bullet_color=col, gap=11, line_spacing=1.05)


# =====================================================================
# SLIDE 5 — GORSELLI BILGI KARTI (HERO)
# =====================================================================
def slide_hero():
    s = base_slide("⭐ Görselli Bilgi Kartı — İmza Özellik", accent=TEAL)
    # sol: pipeline
    lx, lw = 0.6, 5.5
    ly = 1.45
    softcard(s, lx, ly, lw, 5.5, CARD)
    txt(s, lx + 0.35, ly + 0.22, lw - 0.7, 0.5, "Görsel nasıl karta gelir?", size=17, color=INDIGO_D, bold=True)
    steps = [
        ("1", "DIP figürleri ayıklar", "PDF sayfasındaki anatomik figürleri Otsu eşikleme + morfoloji ile arka plandan ayırır."),
        ("2", "En yakın figürü önerir", "Latince/anatomik terime en uygun figürleri bulur — birden çok aday sunar."),
        ("3", "Kullanıcı seçer / kırpar", "Doğru görsel seçilir, gerekirse kırpılır."),
        ("4", "Karta kalıcı eklenir", "Görsel kartın ön/arka yüzüne kalıcı olarak yerleşir."),
    ]
    sy = ly + 0.85
    for num, head, body in steps:
        n = rect(s, lx + 0.35, sy, 0.5, 0.5, TEAL, MSO_SHAPE.OVAL)
        tf = n.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name="Calibri"
        txt(s, lx + 1.0, sy - 0.06, lw - 1.4, 0.4, head, size=14.5, color=INK, bold=True)
        txt(s, lx + 1.0, sy + 0.32, lw - 1.4, 0.65, body, size=11.5, color=SUB, line_spacing=1.0)
        sy += 1.12
    # sag: iki mobil ss yan yana (on yuz + arka yuz)
    rx = lx + lw + 0.5
    rw = SW - rx - 0.6
    place_image(s, M("04-gorselli-kart.png"), rx, 1.45, rw/2 - 0.15, 5.0,
                frame=True, caption="Ön yüz — görsel + soru")
    place_image(s, M("05-kart-cevap.png"), rx + rw/2 + 0.15, 1.45, rw/2 - 0.15, 5.0,
                frame=True, caption="Arka yüz — cevap + SM-2 notları")


# =====================================================================
# SLIDE 6 — WEB ARAYUZU
# =====================================================================
def slide_web():
    s = base_slide("Web Arayüzü", accent=INDIGO)
    shots = [
        (W("01-panel.png"), "Panel — istatistikler ve doküman yönetimi"),
        (W("02-uretim.png"), "Üretim sihirbazı — sayfa aralığı & Gemini"),
        (W("03-deste.png"), "Görselli deste — kart listesi"),
    ]
    cw = (SW - 0.6 - 0.6 - 2*0.4) / 3
    cy = 1.7
    ch = 4.4
    for i, (p, cap) in enumerate(shots):
        x = 0.6 + i * (cw + 0.4)
        place_image(s, p, x, cy, cw, ch, frame=True, caption=cap)


# =====================================================================
# SLIDE 7 — MOBIL UYGULAMA
# =====================================================================
def slide_mobile():
    s = base_slide("Mobil Uygulama (Flutter)", accent=RGBColor(0x2A,0x6F,0xDB))
    # iki telefon solda, metin sagda
    place_image(s, M("01-panel.png"), 0.7, 1.5, 3.1, 5.4, frame=True, caption="Panel — istatistik & ilerleme")
    place_image(s, M("02-ayarlar.png"), 4.05, 1.5, 3.1, 5.4, frame=True, caption="Ayarlar — profil & görünüm")
    tx = 7.5
    tw = SW - tx - 0.6
    softcard(s, tx, 1.6, tw, 5.0, CARD)
    txt(s, tx + 0.4, 1.85, tw - 0.8, 0.5, "Fonksiyonel yaklaşımla Flutter", size=18, color=RGBColor(0x2A,0x6F,0xDB), bold=True)
    bullets(s, tx + 0.45, 2.55, tw - 0.85, 3.9, [
        "Immutability — freezed ile değişmez veri modelleri.",
        "Saf fonksiyonlar — yan etkisiz iş mantığı.",
        "Deklaratif UI — durum → arayüz, tahmin edilebilir.",
        "BLoC — olay/durum akışıyla net mimari.",
        "Aynı API, aynı veri: web ile birebir senkron.",
        "Görselli kartlar mobilde de tam destekli.",
    ], size=14.5, color=INK, bullet_color=RGBColor(0x2A,0x6F,0xDB), gap=12, line_spacing=1.1)


# =====================================================================
# SLIDE 8 — SM-2 / FONKSIYONEL PROGRAMLAMA
# =====================================================================
def slide_sm2():
    s = base_slide("Aralıklı Tekrar (SM-2) — Fonksiyonel Programlama", accent=TEAL)
    # sol metin
    lx, lw = 0.6, 6.2
    softcard(s, lx, 1.5, lw, 5.45, CARD)
    txt(s, lx + 0.4, 1.72, lw - 0.8, 0.5, "Saf fonksiyon olarak SM-2", size=18, color=TEAL_D, bold=True)
    txt(s, lx + 0.4, 2.3, lw - 0.8, 0.9,
        "Bir kartın notu (Tekrar/Zor/İyi/Kolay) ile mevcut durumu girer; bir sonraki tekrar tarihini ve kolaylık faktörünü döndürür. Girdi → çıktı; gizli durum yok.",
        size=13.5, color=SUB, line_spacing=1.12)
    bullets(s, lx + 0.45, 3.35, lw - 0.85, 3.4, [
        "Yan etkisiz — aynı girdi her zaman aynı çıktıyı verir.",
        "Unit testli — algoritma izole biçimde doğrulanır.",
        "Tek algoritma, iki dil — Python (API, otorite) ve Dart (mobil).",
        "Deterministik — planlamada sürpriz yok.",
        "Saf çekirdek, kabukta yan etki — fonksiyonel mimarinin özü.",
    ], size=14, color=INK, bullet_color=TEAL, gap=11, line_spacing=1.08)
    # sag: not butonlari gorseli
    rx = lx + lw + 0.45
    rw = SW - rx - 0.6
    place_image(s, M("05-kart-cevap.png"), rx, 1.5, rw, 5.45, frame=True,
                caption="Notlama: Tekrar / Zor / İyi / Kolay → SM-2 sonraki tarihi hesaplar")


# =====================================================================
# SLIDE 9 — OZELLIKLER
# =====================================================================
def slide_features():
    s = base_slide("Öne Çıkan Özellikler", accent=INDIGO)
    feats = [
        ("📤", "Çok formatlı dışa aktarma", "Anki .apkg, PDF, CSV/TSV/JSON/TXT"),
        ("🌙", "Karanlık mod", "Göz yormayan tema, web & mobil"),
        ("❓", "Quiz / Test sistemi", "Otomatik üretilen çoktan seçmeli"),
        ("🖼️", "Toplu otomatik görsel", "Tüm kartlara tek tıkla görsel"),
        ("⌨️", "Terim otomatik tamamlama", "Anatomik terim önerileri"),
        ("🧭", "Onboarding turu", "İlk kullanımda rehberli tanıtım"),
        ("📈", "İlerleme istatistikleri", "14 günlük çalışma grafiği"),
        ("🔄", "Anlık senkron", "Tek hesap, tüm cihazlar"),
    ]
    # 4x2 grid sol tarafta; sag tarafta ekran goruntusu
    gx, gy = 0.6, 1.6
    gw = 7.8
    cols2, rows2 = 2, 4
    cellw = (gw - 0.4) / cols2
    cellh = 1.22
    for i, (ic, t, d) in enumerate(feats):
        r_ = i % rows2
        c_ = i // rows2
        x = gx + c_ * (cellw + 0.4)
        y = gy + r_ * (cellh + 0.12)
        softcard(s, x, y, cellw, cellh, CARD)
        ib = rect(s, x + 0.2, y + 0.26, 0.7, 0.7, RGBColor(0xEE,0xF1,0xFA), MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = ib.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ic; r.font.size = Pt(20); r.font.name="Calibri"
        txt(s, x + 1.05, y + 0.18, cellw - 1.2, 0.45, t, size=13.5, color=INK, bold=True)
        txt(s, x + 1.05, y + 0.62, cellw - 1.2, 0.5, d, size=10.5, color=SUB, line_spacing=0.95)
    # sag: karanlik mod veya quiz ss
    rx = gx + gw + 0.35
    rw = SW - rx - 0.6
    shot = W("07-karanlik.png") if img_exists(W("07-karanlik.png")) else W("05-quiz.png")
    cap = "Karanlık mod — Panel" if "karanlik" in shot else "Quiz oynatıcı"
    place_image(s, shot, rx, 1.7, rw, 3.0, frame=True, caption=cap)
    if img_exists(W("05-quiz.png")) and "karanlik" in shot:
        place_image(s, W("05-quiz.png"), rx, 4.45, rw, 2.45, frame=True, caption="Quiz oynatıcı")


# =====================================================================
# SLIDE 10 — SENKRONIZASYON KANITI
# =====================================================================
def slide_sync():
    s = base_slide("Senkronizasyon (SSOT) Kanıtı", accent=TEAL)
    # web solda, mobil sagda
    place_image(s, W("01-panel.png"), 0.7, 1.7, 7.0, 4.2, frame=True, caption="Web — Panel")
    place_image(s, M("01-panel.png"), 8.3, 1.55, 2.7, 5.3, frame=True, caption="Mobil — Panel")
    # alt vurgu
    cap = rect(s, 0.7, 6.55, 7.0, 0.55, RGBColor(0xE3,0xF4,0xF1), MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 0.8, 6.62, 6.8, 0.42, "Aynı hesap, aynı veri, anında senkron.",
        size=15, color=TEAL_D, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 11 — DERS ESLEMESI
# =====================================================================
def slide_courses():
    s = base_slide("Ders Eşlemesi", accent=INDIGO)
    cw = (SW - 0.6 - 0.6 - 0.5) / 2
    x1 = 0.6; x2 = 0.6 + cw + 0.5
    cy = 1.6; ch = 5.1
    # Web Tasarim
    softcard(s, x1, cy, cw, ch, CARD)
    rect(s, x1, cy, cw, 1.0, INDIGO, MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, x1 + 0.35, cy + 0.26, cw - 0.7, 0.55, "Web Tasarım ve Programlama", size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x1 + 0.4, cy + 1.2, cw - 0.8, 0.5, "Full-stack istemci–sunucu mimarisi", size=14, color=INDIGO_D, bold=True)
    bullets(s, x1 + 0.45, cy + 1.85, cw - 0.85, ch - 2.1, [
        "React tabanlı SPA istemci",
        "FastAPI ile RESTful API",
        "İlişkisel PostgreSQL veritabanı",
        "Asenkron iş kuyruğu / uzun işlemler",
        "Auth, oturum ve yetkilendirme",
    ], size=14.5, color=INK, bullet_color=INDIGO, gap=12, line_spacing=1.1)
    # Fonksiyonel
    softcard(s, x2, cy, cw, ch, CARD)
    rect(s, x2, cy, cw, 1.0, TEAL, MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, x2 + 0.35, cy + 0.26, cw - 0.7, 0.55, "Fonksiyonel Programlama", size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x2 + 0.4, cy + 1.2, cw - 0.8, 0.5, "Flutter ile fonksiyonel ilkeler", size=14, color=TEAL_D, bold=True)
    bullets(s, x2 + 0.45, cy + 1.85, cw - 0.85, ch - 2.1, [
        "Immutability — freezed değişmez modeller",
        "Saf fonksiyonlar — SM-2 algoritması",
        "Deklaratif UI — durum → arayüz",
        "BLoC — olay/durum akışı",
        "Yan etkisiz, test edilebilir çekirdek",
    ], size=14.5, color=INK, bullet_color=TEAL, gap=12, line_spacing=1.1)


# =====================================================================
# SLIDE 12 — KAPANIS
# =====================================================================
def slide_closing():
    s = add_slide()
    bg = rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, INDIGO)
    grad(bg, (0x00,0x6B,0x60),(0x2A,0x32,0x80), angle=30)
    c2 = rect(s, -1.4, 4.6, 4.0, 4.0, None, MSO_SHAPE.OVAL)
    fill(c2, RGBColor(0x12, 0x9C, 0x8C)); c2.line.fill.background()
    lg = rect(s, SW/2 - 0.55, 1.5, 1.1, 1.1, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = lg.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "M"; r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = INDIGO; r.font.name="Calibri"
    txt(s, 1.0, 2.95, SW - 2.0, 0.9, "MedVisual", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 1.5, 3.95, SW - 3.0, 1.0,
        "Görsel hafızayı tıp eğitimine taşıyan tam entegre ekosistem.",
        size=19, color=RGBColor(0xDDF,0xDD,0xF0) if False else RGBColor(0xDD,0xE6,0xFF), align=PP_ALIGN.CENTER, line_spacing=1.15)
    rect(s, SW/2 - 1.7, 5.05, 3.4, 0.06, RGBColor(0x7C,0xE0,0xD4))
    txt(s, 1.0, 5.35, SW - 2.0, 0.6, "Teşekkürler", size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 1.0, 6.05, SW - 2.0, 0.4, "215541014@firat.edu.tr",
        size=14, color=RGBColor(0xBF,0xE9,0xE2), align=PP_ALIGN.CENTER)


# --- Build ---
slide_cover()
slide_problem()
slide_arch()
slide_stack()
slide_hero()
slide_web()
slide_mobile()
slide_sm2()
slide_features()
slide_sync()
slide_courses()
slide_closing()

prs.save(OUT)
print(f"Kaydedildi: {OUT}")
print(f"Slayt sayisi: {len(prs.slides._sldIdLst)}")
